import os
from pptx import Presentation
import logging

def extract_pptx(file_path: str, document_id: str, images_dir: str):
    """
    Extracts text and images from a PPTX file.
    """
    extracted_data = {
        "pages": [], # Slides
        "headings": [],
        "images": [],
        "diagrams": [],
        "tables": []
    }

    try:
        prs = Presentation(file_path)
        
        for i, slide in enumerate(prs.slides, start=1):
            slide_text = []
            
            # Extract slide title if present
            if slide.shapes.title and slide.shapes.title.text:
                extracted_data["headings"].append({
                    "page_number": i,
                    "level": 1,
                    "text": slide.shapes.title.text
                })

            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_text.append(shape.text.strip())
                
                # Extract images
                if hasattr(shape, "image"):
                    try:
                        image = shape.image
                        img_bytes = image.blob
                        img_ext = image.ext
                        img_filename = f"{document_id}_slide{i}_img_{len(extracted_data['images']) + 1}.{img_ext}"
                        img_path = os.path.join(images_dir, img_filename)
                        
                        with open(img_path, "wb") as f:
                            f.write(img_bytes)
                            
                        extracted_data["images"].append({
                            "page_number": i,
                            "image_path": img_path
                        })
                    except Exception as e:
                        logging.warning(f"Failed to extract PPTX image on slide {i}: {str(e)}")
                        
            extracted_data["pages"].append({
                "page_number": i,
                "text_content": "\n".join(slide_text)
            })

    except Exception as e:
        logging.error(f"Error processing PPTX {file_path}: {str(e)}")
        raise e

    return extracted_data
